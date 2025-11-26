"""
多模态文档处理模块 - 重构版本
支持PDF、Word、PPT、Excel、图片等多种格式的统一处理

主要改进：
- 模块化设计，职责分离
- 统一资源管理和错误处理
- 配置化参数管理
- 更好的类型安全
- 性能优化
"""

import os
import tempfile
import logging
from typing import List, Dict, Any, Optional, Union, Protocol, runtime_checkable
from pathlib import Path
from dataclasses import dataclass, field
from contextlib import contextmanager
from abc import ABC, abstractmethod

import pandas as pd
from langchain.docstore.document import Document

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 延迟导入依赖项
_DEPENDENCIES = {
    'pptx': {'module': 'pptx', 'classes': ['Presentation'], 'error': "python-pptx未安装，PPT处理功能不可用"},
    'docx': {'module': 'docx', 'classes': ['Document'], 'error': "python-docx未安装，Word文档处理功能不可用"},
    'openpyxl': {'module': 'openpyxl', 'error': "openpyxl未安装，Excel处理功能不可用"},
    'tabula': {'module': 'tabula', 'error': "tabula-py未安装，PDF表格提取功能不可用"},
    'camelot': {'module': 'camelot', 'error': "camelot-py未安装，高级PDF表格处理功能不可用"},
    'fitz': {'module': 'fitz', 'error': "PyMuPDF未安装，PDF处理功能不可用"}
}

def _check_dependency(dep_name: str) -> bool:
    """检查依赖项是否可用"""
    try:
        __import__(_DEPENDENCIES[dep_name]['module'])
        return True
    except ImportError:
        logger.warning(_DEPENDENCIES[dep_name]['error'])
        return False

# 依赖项可用性检查
HAS_PPTX = _check_dependency('pptx')
HAS_DOCX = _check_dependency('docx')
HAS_EXCEL = _check_dependency('openpyxl')
HAS_TABULA = _check_dependency('tabula')
HAS_CAMELOT = _check_dependency('camelot')
HAS_PYMUPDF = _check_dependency('fitz')

from .qwen_ocr import QwenOCR


@dataclass
class ProcessingConfig:
    """多模态处理配置"""
    # PDF处理配置
    max_pdf_pages: int = 1000
    max_images_per_page: int = 20
    pdf_table_accuracy_threshold: float = 0.5
    progress_interval: int = 10
    
    # 文本分块配置
    chunk_size: int = 200
    chunk_overlap: int = 0
    
    # OCR配置
    ocr_timeout: float = 30.0
    
    # 临时文件配置
    temp_file_suffix: str = '.png'
    cleanup_temp_files: bool = True
    
    # 支持的文件格式
    image_formats: tuple = ('.jpg', '.jpeg', '.png', '.bmp', '.tiff')
    word_formats: tuple = ('.docx', '.doc')
    ppt_formats: tuple = ('.pptx', '.ppt')
    excel_formats: tuple = ('.xlsx', '.xls')
    pdf_formats: tuple = ('.pdf',)


@dataclass
class ProcessingResult:
    """处理结果统一格式"""
    success: bool
    content: Any = None
    error: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)
    
    @classmethod
    def success_result(cls, content: Any, metadata: Dict[str, Any] = None) -> 'ProcessingResult':
        return cls(success=True, content=content, metadata=metadata or {})
    
    @classmethod
    def error_result(cls, error: str, metadata: Dict[str, Any] = None) -> 'ProcessingResult':
        return cls(success=False, error=error, metadata=metadata or {})


@contextmanager
def managed_temp_file(suffix: str = '.png', cleanup: bool = True):
    """安全的临时文件管理器"""
    temp_file = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    try:
        yield temp_file
    finally:
        temp_file.close()
        if cleanup:
            try:
                os.unlink(temp_file.name)
            except OSError:
                logger.warning(f"无法删除临时文件: {temp_file.name}")


@runtime_checkable
class DocumentProcessor(Protocol):
    """文档处理器协议"""
    
    def can_process(self, file_path: str) -> bool:
        """检查是否可以处理指定文件"""
        ...
    
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        """处理文档文件"""
        ...


class BaseProcessor(ABC):
    """基础处理器抽象类"""
    
    def __init__(self, config: ProcessingConfig):
        self.config = config
    
    @abstractmethod
    def can_process(self, file_path: str) -> bool:
        """检查是否可以处理指定文件"""
        pass
    
    @abstractmethod
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        """处理文档文件"""
        pass
    
    def _get_file_extension(self, file_path: str) -> str:
        """获取文件扩展名"""
        return Path(file_path).suffix.lower()
    
    def _safe_process(self, func, *args, error_msg: str = "处理失败", **kwargs) -> ProcessingResult:
        """安全执行处理函数"""
        try:
            result = func(*args, **kwargs)
            return ProcessingResult.success_result(result)
        except Exception as e:
            logger.error(f"{error_msg}: {str(e)}")
            return ProcessingResult.error_result(f"{error_msg}: {str(e)}")


class ImageProcessor(BaseProcessor):
    """图片处理器"""
    
    def __init__(self, config: ProcessingConfig, ocr_reader: Optional[QwenOCR] = None):
        super().__init__(config)
        self.ocr_reader = ocr_reader
    
    def can_process(self, file_path: str) -> bool:
        return self._get_file_extension(file_path) in self.config.image_formats
    
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        if not self.ocr_reader:
            return ProcessingResult.error_result("OCR功能不可用")
        
        def _extract_text():
            return self.ocr_reader.extract_text(file_path)
        
        return self._safe_process(_extract_text, error_msg="图片文字提取失败")


class PDFProcessor(BaseProcessor):
    """PDF处理器"""
    
    def __init__(self, config: ProcessingConfig, ocr_reader: Optional[QwenOCR] = None):
        super().__init__(config)
        self.ocr_reader = ocr_reader
    
    def can_process(self, file_path: str) -> bool:
        return self._get_file_extension(file_path) in self.config.pdf_formats and HAS_PYMUPDF
    
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        def _process_pdf():
            import fitz
            
            pdf_data = {
                "text_content": [],
                "tables": [],
                "images": []
            }
            
            with fitz.open(file_path) as pdf_document:
                total_pages = len(pdf_document)
                max_pages = min(total_pages, config.max_pdf_pages)
                
                if total_pages > max_pages:
                    logger.warning(f"PDF页数过多({total_pages}页)，仅处理前{max_pages}页")
                
                # 提取文本
                for page_num in range(max_pages):
                    try:
                        page = pdf_document[page_num]
                        pdf_data["text_content"].append({
                            "page_number": page_num + 1,
                            "text": page.get_text()
                        })
                    except Exception as e:
                        logger.warning(f"页面 {page_num+1} 文本提取失败: {e}")
                
                # 提取图片
                if self.ocr_reader:
                    pdf_data["images"] = self._extract_images_from_pdf(pdf_document, max_pages, config)
                
                # 提取表格
                pdf_data["tables"] = self._extract_tables_from_pdf(file_path)
            
            return [pdf_data]
        
        return self._safe_process(_process_pdf, error_msg="PDF处理失败")
    
    def _extract_images_from_pdf(self, pdf_document, max_pages: int, config: ProcessingConfig) -> List[Dict[str, Any]]:
        """从PDF提取图片并OCR"""
        extracted_images = []
        
        for page_num in range(max_pages):
            if page_num % config.progress_interval == 0:
                logger.info(f"正在处理第 {page_num+1}/{max_pages} 页...")
            
            try:
                page = pdf_document[page_num]
                image_list = page.get_images(full=True)
                
                if not image_list:
                    continue
                
                # 限制每页图片数量
                if len(image_list) > config.max_images_per_page:
                    logger.warning(f"第{page_num+1}页图片过多({len(image_list)}张)，仅处理前{config.max_images_per_page}张")
                    image_list = image_list[:config.max_images_per_page]
                
                for img_index, img in enumerate(image_list):
                    image_data = self._extract_single_image(pdf_document, img, page_num, img_index, config)
                    if image_data:
                        extracted_images.append(image_data)
                        
            except Exception as e:
                logger.warning(f"页面 {page_num+1} 处理失败: {e}")
        
        return extracted_images
    
    def _extract_single_image(self, pdf_document, img, page_num: int, img_index: int, config: ProcessingConfig) -> Optional[Dict[str, Any]]:
        """提取单个图片并进行OCR"""
        try:
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            
            if not base_image or not base_image.get("image"):
                return None
            
            with managed_temp_file(config.temp_file_suffix, config.cleanup_temp_files) as temp_file:
                temp_file.write(base_image["image"])
                temp_file.flush()
                
                image_text = self.ocr_reader.extract_text(temp_file.name)
                
                if image_text and image_text.strip() and not image_text.startswith("❌"):
                    return {
                        "page_number": page_num + 1,
                        "image_index": img_index + 1,
                        "filename": f"pdf_page_{page_num+1}_img_{img_index+1}",
                        "extracted_text": image_text.strip(),
                        "image_format": base_image.get("ext", "png"),
                        "image_size": (base_image.get("width", 0), base_image.get("height", 0))
                    }
        except Exception as e:
            logger.warning(f"页面 {page_num+1} 图片 {img_index+1} 处理失败: {e}")
        
        return None
    
    def _extract_tables_from_pdf(self, pdf_path: str) -> List[pd.DataFrame]:
        """从PDF提取表格"""
        tables = []
        
        # 优先使用Camelot
        if HAS_CAMELOT:
            try:
                import camelot
                camelot_tables = camelot.read_pdf(pdf_path, flavor="lattice")
                for table in camelot_tables:
                    if table.parsing_report['accuracy'] > self.config.pdf_table_accuracy_threshold * 100:
                        tables.append(table.df)
                logger.info(f"Camelot提取到 {len(tables)} 个高质量表格")
                return tables
            except Exception as e:
                logger.warning(f"Camelot表格提取失败: {e}")
        
        # 备选方案：使用Tabula
        if HAS_TABULA:
            try:
                import tabula
                tabula_tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
                tables.extend(tabula_tables)
                logger.info(f"Tabula提取到 {len(tables)} 个表格")
            except Exception as e:
                logger.warning(f"Tabula表格提取失败: {e}")
        
        return tables
    

class WordProcessor(BaseProcessor):
    """Word文档处理器"""
    
    def __init__(self, config: ProcessingConfig, ocr_reader: Optional[QwenOCR] = None):
        super().__init__(config)
        self.ocr_reader = ocr_reader
    
    def can_process(self, file_path: str) -> bool:
        return self._get_file_extension(file_path) in self.config.word_formats and HAS_DOCX
    
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        def _process_word():
            # 尝试使用python-docx
            try:
                from docx import Document as DocxDocument
                doc = DocxDocument(file_path)
                return self._process_with_docx(doc)
            except Exception as docx_error:
                logger.warning(f"python-docx处理失败: {docx_error}")
                
                # 如果有OCR功能，尝试使用OCR处理
                if self.ocr_reader:
                    try:
                        logger.info("尝试使用OCR处理Word文档...")
                        text = self.ocr_reader.extract_text(file_path)
                        if text:
                            return [{
                                "paragraphs": [text],
                                "tables": [],
                                "images": []
                            }]
                    except Exception as ocr_error:
                        logger.warning(f"OCR处理失败: {ocr_error}")
                
                # 如果都失败了，尝试使用通用文档加载器
                try:
                    from langchain_community.document_loaders import UnstructuredFileLoader
                    logger.info("尝试使用通用文档加载器...")
                    loader = UnstructuredFileLoader(file_path)
                    docs = loader.load()
                    return [{
                        "paragraphs": [doc.page_content for doc in docs],
                        "tables": [],
                        "images": []
                    }]
                except Exception as unstructured_error:
                    raise Exception(f"所有处理方法都失败: docx错误: {docx_error}, OCR错误: {ocr_error if self.ocr_reader else 'OCR未启用'}, 通用加载器错误: {unstructured_error}")
        
        return self._safe_process(_process_word, error_msg="Word文档处理失败")
    
    def _process_with_docx(self, doc) -> List[Dict[str, Any]]:
        """使用python-docx处理文档"""
        word_data = {
            "paragraphs": [],
            "tables": [],
            "images": []
        }
        
        # 提取段落
        word_data["paragraphs"] = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        
        # 提取表格
        word_data["tables"] = self._extract_word_tables(doc)
        
        # 提取图片
        if self.ocr_reader:
            word_data["images"] = self._extract_word_images(doc, self.config)
        
        return [word_data]
    
    def _extract_word_tables(self, doc) -> List[Dict[str, Any]]:
        """提取Word表格"""
        tables_data = []
        
        for i, table in enumerate(doc.tables):
            try:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                if table_data and len(table_data) > 1:
                    df = pd.DataFrame(table_data[1:], columns=table_data[0])
                    tables_data.append({
                            "table_index": i,
                            "data": df,
                            "summary": f"表格 {i+1} 包含 {len(df)} 行 {len(df.columns)} 列"
                        })
            except Exception as e:
                logger.warning(f"Word表格 {i+1} 处理失败: {e}")
        
        return tables_data
    
    def _extract_word_images(self, doc, config: ProcessingConfig) -> List[Dict[str, Any]]:
        """提取Word图片"""
        images_data = []
        
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_part = rel.target_part
                    
                        with managed_temp_file(config.temp_file_suffix, config.cleanup_temp_files) as temp_file:
                            temp_file.write(image_part.blob)
                            temp_file.flush()
                            
                            image_text = self.ocr_reader.extract_text(temp_file.name)
                            images_data.append({
                                "filename": f"word_image_{len(images_data)+1}",
                                "extracted_text": image_text
                            })
                    except Exception as e:
                        logger.warning(f"Word图片处理失败: {e}")
        except Exception as e:
            logger.warning(f"Word图片提取过程出错: {e}")
        
        return images_data


class PPTProcessor(BaseProcessor):
    """PPT处理器"""
    
    def __init__(self, config: ProcessingConfig, ocr_reader: Optional[QwenOCR] = None):
        super().__init__(config)
        self.ocr_reader = ocr_reader
    
    def can_process(self, file_path: str) -> bool:
        return self._get_file_extension(file_path) in self.config.ppt_formats and HAS_PPTX
    
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        def _process_ppt():
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            slides_data = []
            
            for i, slide in enumerate(presentation.slides):
                slide_content = {
                    "slide_number": i + 1,
                    "title": "",
                    "text_content": [],
                    "images": [],
                    "tables": []
                }
                
                self._extract_slide_content(slide, slide_content, i, config)
                slides_data.append(slide_content)
            
            return slides_data
        
        return self._safe_process(_process_ppt, error_msg="PPT处理失败")
    
    def _extract_slide_content(self, slide, slide_content: Dict, slide_index: int, config: ProcessingConfig):
        """提取幻灯片内容"""
        for shape in slide.shapes:
            # 提取文本
            if hasattr(shape, "text") and shape.text.strip():
                if shape.shape_type == 1:  # 文本框
                    if not slide_content["title"] and len(shape.text) < 100:
                        slide_content["title"] = shape.text.strip()
                    else:
                        slide_content["text_content"].append(shape.text.strip())
            
            # 提取图片
            if shape.shape_type == 13 and self.ocr_reader:  # 图片
                image_data = self._extract_ppt_image(shape, slide_index, len(slide_content["images"]), config)
                if image_data:
                    slide_content["images"].append(image_data)
            
            # 提取表格
            if shape.has_table:
                table_data = self._extract_ppt_table(shape, slide_index)
                if table_data:
                    slide_content["tables"].append(table_data)
    
    def _extract_ppt_image(self, shape, slide_index: int, img_count: int, config: ProcessingConfig) -> Optional[Dict[str, Any]]:
        """提取PPT图片"""
        try:
            image = shape.image
            
            with managed_temp_file(config.temp_file_suffix, config.cleanup_temp_files) as temp_file:
                temp_file.write(image.blob)
                temp_file.flush()
                
                image_text = self.ocr_reader.extract_text(temp_file.name)
                return {
                    "filename": f"slide_{slide_index+1}_image_{img_count+1}",
                    "extracted_text": image_text
                }
        except Exception as e:
            logger.warning(f"幻灯片 {slide_index+1} 图片处理失败: {e}")
        
        return None
    
    def _extract_ppt_table(self, shape, slide_index: int) -> Optional[Dict[str, Any]]:
        """提取PPT表格"""
        try:
            table = shape.table
            table_data = []
            
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data and len(table_data) > 1:
                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                return {
                    "data": df,
                    "summary": f"表格包含 {len(df)} 行 {len(df.columns)} 列"
                }
        except Exception as e:
            logger.warning(f"幻灯片 {slide_index+1} 表格处理失败: {e}")
        
        return None


class ExcelProcessor(BaseProcessor):
    """Excel处理器"""
    
    def can_process(self, file_path: str) -> bool:
        return self._get_file_extension(file_path) in self.config.excel_formats and HAS_EXCEL
    
    def process(self, file_path: str, config: ProcessingConfig) -> ProcessingResult:
        def _process_excel():
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path)
            sheets_data = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                data = []
                
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append(row)
                
                if data:
                    df = pd.DataFrame(data[1:], columns=data[0] if data else None)
                    sheets_data.append({
                        "sheet_name": sheet_name,
                        "data": df,
                        "text_summary": f"工作表 '{sheet_name}' 包含 {len(df)} 行数据"
                    })
            
            return sheets_data
        
        return self._safe_process(_process_excel, error_msg="Excel处理失败")


class MultimodalProcessor:
    """
    多模态文档处理器 - 重构版本
    
    主要改进：
    - 使用策略模式管理不同文档类型的处理器
    - 统一的配置管理
    - 更好的错误处理和资源管理
    - 类型安全的接口设计
    """
    
    def __init__(self, api_key: Optional[str] = None, config: Optional[ProcessingConfig] = None):
        self.config = config or ProcessingConfig()
        self.ocr_reader = self._init_ocr(api_key)
        self.processors = self._init_processors()
    
    def _init_ocr(self, api_key: Optional[str]) -> Optional[QwenOCR]:
        """初始化OCR引擎"""
        try:
            ocr_reader = QwenOCR(api_key)
            logger.info("OCR引擎初始化成功 (使用 Qwen2.5-VL)")
            return ocr_reader
        except Exception as e:
            logger.warning(f"OCR引擎初始化失败: {e}")
            return None
    
    def _init_processors(self) -> List[DocumentProcessor]:
        """初始化文档处理器"""
        return [
            ImageProcessor(self.config, self.ocr_reader),
            PDFProcessor(self.config, self.ocr_reader),
            WordProcessor(self.config, self.ocr_reader),
            PPTProcessor(self.config, self.ocr_reader),
            ExcelProcessor(self.config)
        ]
    
    def get_processor_for_file(self, file_path: str) -> Optional[DocumentProcessor]:
        """获取适合的文档处理器"""
        for processor in self.processors:
            if processor.can_process(file_path):
                return processor
        return None
    
    def process_file(self, file_path: str) -> ProcessingResult:
        """处理文档文件"""
        processor = self.get_processor_for_file(file_path)
        
        if not processor:
            return ProcessingResult.error_result(f"不支持的文件格式: {Path(file_path).suffix}")
        
        logger.info(f"使用 {processor.__class__.__name__} 处理文件: {file_path}")
        return processor.process(file_path, self.config)
    
    def extract_text_from_image(self, image_path: str) -> str:
        """OCR识别图片中的文字 - 保持向后兼容"""
        if not self.ocr_reader:
            return "⚠️  OCR功能不可用"
        
        try:
            return self.ocr_reader.extract_text(image_path)
        except Exception as e:
            return f"❌ 图片文字提取失败: {str(e)}"


def get_multimodal_supported_formats() -> Dict[str, str]:
    """返回支持的多模态格式"""
    config = ProcessingConfig()
    formats = {}
    
    # 图片格式
    if HAS_PYMUPDF:
        for ext in config.image_formats:
            formats[ext] = '图片文件 (OCR文字识别)'
    
    # Word格式
    if HAS_DOCX:
        for ext in config.word_formats:
            formats[ext] = 'Word文档 (文本+图片+表格)'
    
    # PPT格式
    if HAS_PPTX:
        for ext in config.ppt_formats:
            formats[ext] = 'PowerPoint演示文稿 (文本+图片+表格)'
    
    # Excel格式
    if HAS_EXCEL:
        for ext in config.excel_formats:
            formats[ext] = 'Excel工作簿 (表格数据)'
    
    # PDF格式
    pdf_features = []
    if HAS_TABULA or HAS_CAMELOT:
        pdf_features.append("表格提取")
    if HAS_PYMUPDF:
        pdf_features.append("图片OCR")
    
    if pdf_features:
        formats['.pdf'] = f'PDF文档 (文本+{"+".join(pdf_features)})'
    
    return formats


def create_multimodal_documents(file_path: str, processor: MultimodalProcessor) -> List[Document]:
    """
    创建多模态文档对象
    
    Args:
        file_path: 文件路径
        processor: 多模态处理器实例
        
    Returns:
        处理后的文档列表
    """
    result = processor.process_file(file_path)
    
    if not result.success:
        logger.error(f"文档处理失败: {result.error}")
        return []
    
    documents = []
    file_extension = Path(file_path).suffix.lower()
    
    try:
        if file_extension in processor.config.image_formats:
            # 图片文件
            documents.append(Document(
                page_content=f"图片文字识别结果:\n{result.content}",
                metadata={"source": file_path, "type": "image_ocr"}
            ))
        
        elif file_extension in processor.config.word_formats:
            # Word文件
            documents.extend(_create_word_documents(result.content, file_path))
        
        elif file_extension in processor.config.ppt_formats:
            # PPT文件
            documents.extend(_create_ppt_documents(result.content, file_path))
        
        elif file_extension in processor.config.excel_formats:
            # Excel文件
            documents.extend(_create_excel_documents(result.content, file_path))
        
        elif file_extension in processor.config.pdf_formats:
            # PDF文件
            documents.extend(_create_pdf_documents(result.content, file_path))
        
        logger.info(f"多模态处理完成，生成 {len(documents)} 个文档片段")
        return documents
        
    except Exception as e:
        logger.error(f"创建文档对象失败: {str(e)}")
        return []


def _create_word_documents(word_data_list: List[Dict], file_path: str) -> List[Document]:
    """创建Word文档对象"""
    documents = []
            
    for word_data in word_data_list:
        if "error" in word_data:
            continue
        
        content_parts = []
        
        # 添加段落文本
        if word_data.get("paragraphs"):
            content_parts.append("文本内容:\n" + "\n".join(word_data["paragraphs"]))
        
        # 添加图片OCR内容
        if word_data.get("images"):
            image_texts = []
            for i, img in enumerate(word_data["images"], 1):
                extracted_text = str(img.get("extracted_text", "")).strip()
                if extracted_text:
                    image_texts.append(f"第{i}张图片 ({img['filename']})：\n{extracted_text}")
            if image_texts:
                content_parts.append("图片文字内容：\n" + "\n\n".join(image_texts))
        
        # 添加表格内容
        if word_data.get("tables"):
            table_texts = []
            for table in word_data["tables"]:
                table_texts.append(f"第{table['table_index']+1}个表格:\n{table['summary']}\n{'-'*50}\n{table['data'].to_string()}")
            if table_texts:
                content_parts.append("表格内容:\n" + "\n\n".join(table_texts))
                
        if content_parts:
            documents.append(Document(
                page_content="\n\n".join(content_parts),
                metadata={
                    "source": file_path,
                    "type": "word_document",
                    "has_images": len(word_data.get("images", [])) > 0,
                    "has_tables": len(word_data.get("tables", [])) > 0,
                    "paragraphs_count": len(word_data.get("paragraphs", []))
                }
            ))
    
    return documents


def _create_ppt_documents(slides_data: List[Dict], file_path: str) -> List[Document]:
    """创建PPT文档对象"""
    documents = []
            
    for slide_data in slides_data:
        if "error" in slide_data:
            continue
        
        content_parts = []
        
        if slide_data.get("title"):
            content_parts.append(f"标题: {slide_data['title']}")
        
        if slide_data.get("text_content"):
            content_parts.append("文本内容:\n" + "\n".join(slide_data["text_content"]))
        
        if slide_data.get("images"):
            for img in slide_data["images"]:
                extracted_text = str(img.get("extracted_text", "")).strip()
                if extracted_text:
                    content_parts.append(f"图片文字: {extracted_text}")
        
        if slide_data.get("tables"):
            for table in slide_data["tables"]:
                content_parts.append(f"表格数据: {table['summary']}\n{table['data'].to_string()}")
                
        if content_parts:
            documents.append(Document(
                page_content="\n\n".join(content_parts),
                metadata={
                    "source": file_path, 
                    "type": "ppt_slide",
                    "slide_number": slide_data["slide_number"]
                }
            ))
        
    return documents


def _create_excel_documents(sheets_data: List[Dict], file_path: str) -> List[Document]:
    """创建Excel文档对象"""
    documents = []
            
    for sheet_data in sheets_data:
        if "error" in sheet_data:
            continue
        
        df = sheet_data["data"]
        content = f"工作表名称: {sheet_data['sheet_name']}\n"
        content += f"数据概览: {sheet_data['text_summary']}\n"
        content += f"表格内容:\n{df.to_string()}"
        
        documents.append(Document(
            page_content=content,
            metadata={
                "source": file_path,
                "type": "excel_sheet", 
                "sheet_name": sheet_data["sheet_name"]
            }
        ))
        
    return documents


def _create_pdf_documents(pdf_data_list: List[Dict], file_path: str) -> List[Document]:
    """创建PDF文档对象"""
    documents = []
                
    for pdf_data in pdf_data_list:
        if "error" in pdf_data:
            continue
        
        content_parts = []
        
        # 添加文本内容
        if pdf_data.get("text_content"):
            text_parts = []
            for text_page in pdf_data["text_content"]:
                if isinstance(text_page, dict) and "page_number" in text_page:
                    text_parts.append(f"第{text_page['page_number']}页:\n{text_page['text']}")
            if text_parts:
                content_parts.append("文本内容:\n" + "\n\n".join(text_parts))
        
        # 添加表格内容
        if pdf_data.get("tables"):
            table_texts = []
            for i, table in enumerate(pdf_data["tables"]):
                if hasattr(table, 'to_string'):
                    table_texts.append(f"表格 {i+1}: {len(table)} 行 {len(table.columns)} 列\n{table.to_string()}")
            if table_texts:
                content_parts.append("表格内容:\n" + "\n\n".join(table_texts))
        
        # 添加图片OCR内容
        if pdf_data.get("images"):
            image_texts = []
            for img in pdf_data["images"]:
                if isinstance(img, dict) and "page_number" in img:
                    image_texts.append(f"第{img['page_number']}页图片{img.get('image_index', 0)}: {img.get('extracted_text', '')}")
            if image_texts:
                content_parts.append("图片文字内容:\n" + "\n".join(image_texts))
        
        if content_parts:
            documents.append(Document(
                page_content="\n\n".join(content_parts),
                metadata={
                    "source": file_path,
                    "type": "pdf_multimodal",
                    "has_text": len(pdf_data.get("text_content", [])) > 0,
                    "has_tables": len(pdf_data.get("tables", [])) > 0,
                    "has_images": len(pdf_data.get("images", [])) > 0,
                    "pages_count": len(pdf_data.get("text_content", []))
                }
            ))
        
    return documents


# 全局处理器实例
_global_processor: Optional[MultimodalProcessor] = None

def get_multimodal_processor() -> MultimodalProcessor:
    """获取全局多模态处理器实例"""
    global _global_processor
    if _global_processor is None:
        _global_processor = MultimodalProcessor()
    return _global_processor 